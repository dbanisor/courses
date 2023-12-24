import collections.abc
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt, Cm
import pptx.enum.shapes
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.shapes.picture import Picture
from pptx.enum.dml import MSO_THEME_COLOR_INDEX

prs = Presentation()
prs.slide_width = Inches(8.2)                   # cm = 20.8 cm
# print(prs.slide_width.cm)
prs.slide_height = Inches(19.25)                # cm = 48.9 cm
# print(prs.slide_height.cm)
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)


shapes = slide.shapes
left = Cm(5)
top = Cm(5)
width = height = Cm(1.0)
shape = shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
fill = shape.fill
fill.solid()
fill.fore_color.theme_color = MSO_THEME_COLOR_INDEX.ACCENT_1
fill.fore_color.brightness = -0.25
line = shape.line
line.color.rgb = RGBColor(255, 0, 0)
line.width = Pt(2.5)


header = Picture('headers/1000_F_268488616_wcoB2JnGbOD2u3bpn2GPmu0KJQ4Ah66T.jpg', prs)

prs.save('test.pptx')